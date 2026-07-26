"""Parses PPTX bytes into slide-level `DocumentSegment`s.

Extracts both free-text shapes and table cell content, since order details
in the fixture set appear in tables as often as in text boxes.
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from src.parsers import ParserError
from src.schema import DocumentSegment, ParsedDocument


def _shape_text(shape) -> list[str]:
    if shape.has_table:
        return [
            cell.text.strip()
            for row in shape.table.rows
            for cell in row.cells
            if cell.text.strip()
        ]
    if shape.has_text_frame and shape.text_frame.text.strip():
        return [shape.text_frame.text.strip()]
    return []


def parse_pptx(raw_bytes: bytes) -> ParsedDocument:
    """Parses PPTX bytes into one `DocumentSegment` per slide.

    Args:
        raw_bytes: Full contents of a `.pptx` file.

    Returns:
        A `ParsedDocument` with `locator` values of the form `slide:<n>`
        (1-indexed), text joined from all text boxes and table cells on
        that slide.

    Raises:
        ParserError: If the bytes are not a valid PPTX package.
    """
    try:
        presentation = Presentation(BytesIO(raw_bytes))
    except Exception as exc:
        raise ParserError(f"failed to open PPTX: {exc}") from exc

    segments = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines = [line for shape in slide.shapes for line in _shape_text(shape)]
        segments.append(DocumentSegment(locator=f"slide:{slide_index}", text="\n".join(lines)))

    if not segments:
        raise ParserError("PPTX has no slides")

    return ParsedDocument(source_type="pptx", segments=segments)
